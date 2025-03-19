import streamlit as st
import openpyxl
import pyperclip
import os
import shutil
import pandas as pd
import io
import fitz  # PyMuPDF
import re
from io import BytesIO
from docx import Document
from pptx import Presentation

# 페이지 설정
st.set_page_config(page_title="엑셀 도구 모음", layout="centered")

# 사이드바 메뉴
st.sidebar.title("엑셀 도구 모음")
page = st.sidebar.radio(" ", ("엑셀 데이터 복사", "엑셀 시트 분할", "단어수 카운터(파일)", "단어수 카운터(웹)", "월간 보고 데이터"))

# 1. 엑셀 데이터 복사
if page == "엑셀 데이터 복사":
    st.title('📄엑셀 데이터 복사')
    st.write(":rainbow[지정된 키워드 바로 아래 행부터 전체 내용이 클립보드에 복사됩니다.]")

    # 파일 업로드
    uploaded_file = st.file_uploader("엑셀 파일 업로드", type=["xlsx"])
   
    default_keywords = ["중간_CNS", "zh-hans", "CNS", "zh_CN", "Simplified Chinese", "CNS (중국어 간체)"]  # 기본 키워드
    # 키워드 입력 (사용자가 수정할 수 있음)
    keywords_input = st.text_area("찾을 키워드(언어열 이름)", value=", ".join(default_keywords))

    if uploaded_file:
        file_name = uploaded_file.name
        directory_path = os.getcwd()  # 현재 작업 디렉토리로 설정
            
        # 키워드 처리: 쉼표로 구분된 키워드 리스트로 변환
        keywords = [keyword.strip() for keyword in keywords_input.split(',')]

        # 파일명에 확장자 .xlsx가 없다면 자동으로 추가
        if not file_name.lower().endswith(".xlsx"):
            file_name += ".xlsx"
                
        # 다운로드 폴더 경로 설정
        downloads_folder = os.path.join(os.path.expanduser("~"), "Downloads")
        file_path = os.path.join(downloads_folder, file_name)

        if os.path.exists(file_path):
         
            # 엑셀 파일 열기
                wb = openpyxl.load_workbook(file_path, data_only=True)
                ws = wb.active  # 첫 번째 시트 활성화

                # 특정 키워드가 있는 행과 열 찾기
                target_row = None
                target_column = None

                for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=ws.max_row, values_only=True), start=1):
                    for col_idx, cell_value in enumerate(row, start=1):
                        if cell_value in keywords:
                            target_row = row_idx
                            target_column = col_idx
                            break
                    if target_row:
                        break  # 첫 번째 일치하는 키워드만 찾음

                # 키워드가 발견되지 않으면 종료
                if target_row is None:
                    st.error("❌ 키워드를 찾을 수 없습니다.")
                else:
                    # 키워드 행 바로 아래부터 끝까지 해당 열의 데이터 가져오기
                    values = [
                        str(ws.cell(row=i, column=target_column).value).replace("\n", "\r\n") 
                        if ws.cell(row=i, column=target_column).value is not None else ""
                        for i in range(target_row + 1, ws.max_row + 1)
                    ]
                    
                    # 엑셀에서 붙여넣을 때 한 셀 안에 줄바꿈이 유지되도록 " "로 감싸기
                    if values:
                        formatted_text = "\r\n".join(f'"{value}"' for value in values)  # 각 셀 값을 " "로 감싸기
                        pyperclip.copy(formatted_text)  # 클립보드에 복사   
                    
                        # 성공 메시지와 formatted_text를 다른 영역에 표시
                        st.success("✅ 클립보드에 복사 완료!")  # 성공 메시지
                        st.text_area("복사된 내용", formatted_text, height=200)  # 복사된 내용 표시 
                    else:
                        st.warning("⚠️ 복사할 데이터가 없습니다.")

                # 워크북 닫기
                wb.close()
        else:
            st.error("❌ 파일이 존재하지 않습니다.")
            directory_path = None        

# 2. 엑셀 시트 분할 페이지
elif page == "엑셀 시트 분할":
    st.title("✂️ 엑셀 시트 분할")
    st.caption("※ 엑셀 파일의 각 시트를 새로운 파일로 분할하여 저장합니다.")

    # 파일 업로드
    uploaded_file = st.file_uploader("엑셀 파일을 업로드하세요 (.xlsx)", type=["xlsx"])
    if uploaded_file:
        file_name = uploaded_file.name
    else:
        file_name = None

    # 실행 버튼 클릭 시 동작
    if file_name and st.button("🚀 실행", use_container_width=True):
        # 업로드된 파일을 메모리에서 읽을 수 있도록 처리
        excel_file = pd.ExcelFile(io.BytesIO(uploaded_file.read()))

        # 파일이 존재하는지 확인
        if not excel_file.sheet_names:
            st.error("❌ 해당 엑셀 파일에 시트가 없습니다.")
        else:
            # 새로운 폴더로 시트 분할하여 저장
            output_folder = os.path.join(os.getcwd(), f"{os.path.splitext(file_name)[0]}_시트분할")

            # 기존 폴더 삭제 후 새로 생성
            if os.path.exists(output_folder):
                shutil.rmtree(output_folder)
            os.makedirs(output_folder)

            # 시트별 저장
            for sheet_name in excel_file.sheet_names:
                # 각 시트의 데이터를 DataFrame으로 가져오기
                df = excel_file.parse(sheet_name)

                # 시트별로 새로운 엑셀 파일로 저장
                output_path = os.path.join(output_folder, f"{sheet_name}.xlsx")
                df.to_excel(output_path, index=False, sheet_name=sheet_name)

            st.success(f"✅ 모든 시트가 '{output_folder}' 폴더에 저장되었습니다.")

# 3. 단어수 카운트
elif page == "단어수 카운터":
    st.title("🔢 단어수 카운터")

    import fitz  # PyMuPDF
    from docx import Document
    from pptx import Presentation
    import openpyxl

    def count_words_in_word(file):
        try:
            doc = Document(file)
            words = 0
            content = []
            for para in doc.paragraphs:
                words += len(para.text.split())
                content.append(para.text[:300] + "\n")
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        words += len(cell.text.split())
                        content.append(cell.text[:300] + "\n")
            return words, "".join(content)
        except Exception as e:
            return 0, f"Word 파일 처리 중 오류: {e}"

    def count_words_in_pptx(file):
        try:
            prs = Presentation(file)
            words = 0
            content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        words += len(shape.text.replace("-", " ").split())
                        content.append(shape.text[:300] + "\n")
            return words, "".join(content)
        except Exception as e:
            return 0, f"PPTX 파일 처리 중 오류: {e}"

    def count_words_in_excel(file):
        try:
            wb = openpyxl.load_workbook(file)
            words = 0
            content = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value:
                            words += len(str(cell.value).split())
                            content.append(str(cell.value)[:300] + "\n")
            return words, "".join(content)
        except Exception as e:
            return 0, f"Excel 파일 처리 중 오류: {e}"

    def count_words_in_pdf(file):
        try:
            doc = fitz.open(file)
            words = 0
            content = []
            for page in doc:
                text = page.get_text()
                words += len(text.split())
                content.append(text[:300] + "\n")
            return words, "".join(content)
        except Exception as e:
            return 0, f"PDF 파일 처리 중 오류: {e}"

    def count_words_in_txt(file):
        try:
            content = file.read().decode("utf-8")
            words = len(content.split())
            return words, content[:300] + "\n"
        except Exception as e:
            return 0, f"TXT 파일 처리 중 오류: {e}"

    st.markdown("<span style='font-size: 20px;'>파일을 업로드하세요. (지원 형식: Word, PPTX, Excel, PDF, TXT)</span>", unsafe_allow_html=True)
    st.caption("※ 여러 언어가 섞인 파일인 경우 단어수가 부정확하게 나옵니다.")
    uploaded_file = st.file_uploader("", type=["pptx", "docx", "xlsx", "pdf", "txt"])

    if uploaded_file is not None:
        file_name = uploaded_file.name.lower()
        word_count = 0
        file_content = ""

        if file_name.endswith(".docx"):
            word_count, file_content = count_words_in_word(uploaded_file)
            file_type_display = "📄 Word"
        elif file_name.endswith(".pptx"):
            word_count, file_content = count_words_in_pptx(uploaded_file)
            file_type_display = "📊 PPTX"
        elif file_name.endswith(".xlsx"):
            word_count, file_content = count_words_in_excel(uploaded_file)
            file_type_display = "📑 Excel"
        elif file_name.endswith(".pdf"):
            word_count, file_content = count_words_in_pdf(uploaded_file)
            file_type_display = "📜 PDF"
        elif file_name.endswith(".txt"):
            word_count, file_content = count_words_in_txt(uploaded_file)
            file_type_display = "📃 TXT"
        else:
            file_type_display = "❌ 지원되지 않는 파일"

        st.markdown(f"### 파일 형식: {file_type_display}")
        st.markdown(f"### 단어수: <span style='color: #4CAF50; font-size: 24px;'>{word_count}</span>", unsafe_allow_html=True)
        st.text_area("파일 내용 미리보기:", file_content, height=200)

# 4. 월간 보고 데이터
elif page == "월간 보고 데이터":
# Streamlit UI 구성
    st.title("📊 Jira CSV 데이터 추출기")
    st.write("CSV 파일을 업로드하면 필요한 데이터만 추출하여 엑셀로 다운로드할 수 있습니다.")

    # "요약"에서 단어수와 기준 언어 추출
    def extract_info_from_summary(요약):
        if pd.isna(요약) or not isinstance(요약, str):  # NaN 체크 및 문자열 확인
            return pd.Series([None, None])
        
        match = re.search(r'\[(\d+)\s*([A-Za-z]+)\]', 요약)  # 단어수와 기준 언어 추출
        if match:
            word_count = match.group(1)  # 숫자 (단어수)
            language = match.group(2)  # 영어 코드 (기준 언어)
            return pd.Series([word_count, language])
        return pd.Series([None, None])



    # 파일 업로드
    uploaded_file = st.file_uploader("📂 CSV 파일을 업로드하세요", type=["csv"])

    if uploaded_file is not None:
        # CSV 읽기
        df = pd.read_csv(uploaded_file)

        # 필요한 열만 선택 (해당 열이 없으면 기본값 사용)
        selected_columns = ['프로젝트 이름', '요약', '기한', '생성일']
        for col in selected_columns:
            if col not in df.columns:
                df[col] = None  # 없는 열은 빈 값 생성

        df_filtered = df[selected_columns].copy()

        # "기한"과 "생성일" 열에서 처음 10개의 문자만 추출
        df_filtered['기한'] = df_filtered['기한'].astype(str).str[:10]
        df_filtered['생성일'] = df_filtered['생성일'].astype(str).str[:10]

        # "요약"에서 단어수와 기준 언어 추출
        df_filtered[['단어수', '기준 언어']] = df_filtered['요약'].apply(extract_info_from_summary)

        # '단어수' 열을 숫자로 변환
        df_filtered['단어수'] = pd.to_numeric(df_filtered['단어수'], errors='coerce')

        # 프로젝트별로 항목 개수 및 단어수 합계 계산
        project_summary_df = df_filtered.groupby('프로젝트 이름').agg(
            요청수=('요약', 'count'),
            단어수_합계=('단어수', 'sum')
        ).reset_index()

        # 요청수 열 기준으로 내림차순 정렬
        project_summary_df = project_summary_df.sort_values(by='요청수', ascending=False)

        # 합계 행 추가
        total_row = project_summary_df[['요청수', '단어수_합계']].sum().to_frame().T
        total_row['프로젝트 이름'] = '합계'
        project_summary_df = pd.concat([project_summary_df, total_row], ignore_index=True)

        # 데이터 미리보기
        st.write("📋 **원본 데이터**")
        st.dataframe(df_filtered.head())  # 상위 5개 데이터 미리 보기

        # 엑셀로 변환
        output = BytesIO()
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            # 원본 데이터 시트
            df_filtered.to_excel(writer, index=False, sheet_name="원본 데이터")
            
            # 프로젝트별 요약 시트
            project_summary_df.to_excel(writer, index=False, sheet_name="프로젝트별 요약")
            
            writer.close()
        output.seek(0)

        # 다운로드 버튼
        st.download_button(
            label="📥 엑셀 파일 다운로드",
            data=output,
            file_name="project_summary.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

# 5. 단어수 카운터
elif page == "단어수 카운터(웹)":
    st.title("🔢 단어수 카운터(웹)")
    st.write("텍스트를 입력하면 띄어쓰기 기준으로 단어 수를 계산합니다.")

    def count_words(text):
        words = text.split()
        return len(words)

    if 'word_count' not in st.session_state:
        st.session_state.word_count = 0

    st.subheader(f"단어 수: {st.session_state.word_count}")

    def update_word_count():
        st.session_state.word_count = count_words(st.session_state.text_input)

    text_input = st.text_area("텍스트 입력", height=200, key='text_input', on_change=update_word_count)
